from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


PRODUCTS = [
    {
        "rank": "01",
        "name": "众安尊享e生系列",
        "insurer": "众安在线",
        "proxy": "公开累计服务用户超7200万，百万医疗老牌产品线",
        "why": ["公开用户规模证据最强", "互联网百万医疗代表性产品", "产品迭代多年，渠道认知度高"],
        "evidence": "S3主流媒体规模证据 + S1公开产品材料",
        "source": "中国证券网/上海证券报：尊享e生9年累计服务用户超7200万",
    },
    {
        "rank": "02",
        "name": "好医保长期医疗旗舰版/2025",
        "insurer": "人保健康 + 蚂蚁保",
        "proxy": "蚂蚁保强渠道 + 人保健康承保 + 长期医疗持续升级",
        "why": ["平台触达强，用户认知度高", "长期医疗形态持续升级", "成人/中老年/少儿产品矩阵覆盖广"],
        "evidence": "S3主流媒体产品升级证据 + S2平台/承保方信号",
        "source": "中国经济网：人保健康联合蚂蚁保推出旗舰版、2025版升级",
    },
    {
        "rank": "03",
        "name": "蓝医保长期医疗/好医好药版",
        "insurer": "太平洋健康险",
        "proxy": "太保健康主推互联网医疗险，2026继续升级",
        "why": ["官方持续升级信号明确", "好医好药版强化院外药与服务", "品牌和渠道曝光度高"],
        "evidence": "S2官方升级公告 + S3主流媒体",
        "source": "太平洋保险官网/新华网：2026蓝医保与好医好药版升级",
    },
    {
        "rank": "04",
        "name": "平安e生保百万医疗2026",
        "insurer": "平安健康险",
        "proxy": "平安健康险热销产品线，e生保系列服务客户规模公开",
        "why": ["平安品牌和健康险渠道优势", "官方产品页持续展示", "系列累计服务与赔付规模公开"],
        "evidence": "S1官方产品页 + S3主流媒体规模证据",
        "source": "平安官网/央广网/新浪财经：e生保2026与系列服务规模",
    },
    {
        "rank": "05",
        "name": "腾讯微保微医保系列",
        "insurer": "腾讯微保平台 + 多承保方",
        "proxy": "平台型医疗险矩阵，公开超1700万用户健康保障服务",
        "why": ["微信生态渠道触达强", "微医保品牌矩阵化", "2026系列升级信号明确"],
        "evidence": "S3公开用户规模证据 + S1条款样本",
        "source": "投资界/中证网/凤凰网/泰康在线：微医保用户规模、升级与条款",
    },
]

SOURCES = [
    "中国证券网/上海证券报：https://www.cnstock.com/commonDetail/15491",
    "中国经济网：https://finance.ce.cn/insurance1/scrollnews/202404/01/t20240401_38955106.shtml",
    "中国经济网：https://finance.ce.cn/insurance1/scrollnews/202506/t20250619_2333272.shtml",
    "太平洋保险官网：https://www.cpic.com.cn/c/2025-12-22/1879211.shtml",
    "新华网：https://app.xinhuanet.com/news/article.html?articleId=1e65036fae09c4e9d2518311ca553cbe",
    "平安官网：https://www.pingan.com/official/new-product/MNg%3D%3D",
    "央广网：https://finance.cnr.cn/zghq/20260601/t20260601_527643499.shtml",
    "新浪财经：https://finance.sina.cn/2025-12-01/detail-infzikpw2840678.d.html",
    "投资界：https://news.pedaily.cn/20230531/59420.shtml",
    "中证网：https://cs.com.cn/bx/202408/t20240816_6432188.html",
    "凤凰网：https://finance.ifeng.com/c/8sOlDKUhUZ7",
    "泰康在线条款：https://policy.taikang.com/tkcms/article/70749b9f1139b82735dfefd5c2ccd8fa.pdf",
]

BG = RGBColor(248, 250, 252)
NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(14, 165, 233)
GRAY = RGBColor(71, 85, 105)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def textbox(slide, x, y, w, h, text, size=14, color=GRAY, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return box


def add_title(slide, title, subtitle=None):
    textbox(slide, 0.55, 0.35, 12.1, 0.65, title, 26, NAVY, True)
    if subtitle:
        textbox(slide, 0.58, 0.95, 11.9, 0.4, subtitle, 11, GRAY)


def card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(226, 232, 240)
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    textbox(slide, x + 0.25, y + 0.2, w - 0.4, 0.35, title, 14, NAVY, True)
    textbox(slide, x + 0.25, y + 0.65, w - 0.4, h - 0.8, body, 11, GRAY)


def build_deck(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    textbox(slide, 0.75, 1.15, 11.8, 0.7, "公开证据口径下的热销医疗险 Top5", 34, NAVY, True)
    textbox(slide, 0.78, 2.05, 11.0, 0.5, "个人商业百万医疗/互联网医疗险市场研究，不构成投保建议", 17, GRAY)
    card(slide, 0.78, 3.0, 3.8, 1.45, "核心口径", "没有公开单品真实销量表；本报告采用公开用户规模、渠道覆盖、官方升级和材料可得性作为热销代理。", BLUE)
    card(slide, 4.9, 3.0, 3.8, 1.45, "输出对象", "尊享e生、好医保、蓝医保、平安e生保、微医保系列。", CYAN)
    card(slide, 9.02, 3.0, 3.8, 1.45, "研究日期", "2026-06-04；资料以后续官方披露为准。", GREEN)
    textbox(slide, 0.78, 6.85, 11.5, 0.3, "提示：Top5为公开证据代理排序，不是内部销量排名。", 10, ORANGE, True)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "方法论：把“卖得最多”转成可复核代理指标")
    items = [
        ("1. 公开规模", "累计服务用户、客户数、赔付规模等公开披露优先。"),
        ("2. 渠道覆盖", "支付宝、微信、平安、太保、众安等渠道或品牌触达。"),
        ("3. 持续在售/升级", "2025-2026 仍有官方升级、产品页或条款更新。"),
        ("4. 官方材料", "有条款、产品页、承保方或官方新闻可回溯。"),
        ("5. 风险控制", "不输出购买建议，只做B端市场研究。"),
    ]
    for i, (title, body) in enumerate(items):
        card(slide, 0.75 + (i % 3) * 4.2, 1.55 + (i // 3) * 2.05, 3.75, 1.45, title, body, [BLUE, CYAN, GREEN, ORANGE, BLUE][i])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Top5 结论矩阵", "按公开证据强度与热销代理信号排序")
    headers = ["排名", "产品/系列", "承保/平台", "为什么像热销", "证据等级"]
    xs = [0.55, 1.35, 4.15, 6.35, 10.65]
    ws = [0.7, 2.6, 2.1, 4.0, 2.0]
    y = 1.35
    for x, w, header in zip(xs, ws, headers):
        textbox(slide, x, y, w, 0.35, header, 10, NAVY, True)
    y += 0.45
    for product in PRODUCTS:
        vals = [product["rank"], product["name"], product["insurer"], product["proxy"], product["evidence"]]
        for x, w, val in zip(xs, ws, vals):
            textbox(slide, x, y, w, 0.78, val, 9, GRAY, product["rank"] == "01")
        y += 0.93

    for product in PRODUCTS:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_title(slide, f"{product['rank']} {product['name']}", product["insurer"])
        textbox(slide, 0.75, 1.45, 2.0, 0.6, product["rank"], 36, BLUE, True)
        card(slide, 2.25, 1.35, 5.0, 1.6, "入选理由", product["proxy"], BLUE)
        card(slide, 7.65, 1.35, 4.85, 1.6, "证据强度", product["evidence"], GREEN)
        card(slide, 0.75, 3.15, 6.1, 2.1, "为什么卖得动", "\n".join([f"• {x}" for x in product["why"]]), CYAN)
        card(slide, 7.2, 3.15, 5.3, 2.1, "关键来源", product["source"], ORANGE)
        textbox(slide, 0.8, 6.55, 11.8, 0.35, "研究提示：这里是市场热度判断，不代表适合任何具体个人购买。", 10, ORANGE, True)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "已固化为可复用 Tool 流程")
    steps = [
        "定义范围：个人商业百万医疗/互联网医疗险，排除基本医保和惠民保",
        "收集候选：官方材料、承保方/平台公告、主流财经媒体、第三方弱信号",
        "证据分级：S1官方条款/页面，S2官方公告，S3主流媒体，S5第三方榜单",
        "代理排序：公开规模 > 渠道覆盖 > 近期升级 > 官方材料可得性",
        "产出：TopN矩阵、单品理由、限制说明、来源附录和PPT",
    ]
    textbox(slide, 0.85, 1.45, 11.7, 4.7, "\n".join([f"{i + 1}. {step}" for i, step in enumerate(steps)]), 17, GRAY)
    textbox(slide, 0.85, 6.55, 11.7, 0.4, "Tool spec: docs/tools/medical-insurance-market-research-tool.md", 11, BLUE, True)

    for part_idx in range(0, len(SOURCES), 6):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_title(slide, "来源附录", "用于验证公开规模、产品升级、官方材料或条款")
        chunk = SOURCES[part_idx : part_idx + 6]
        textbox(slide, 0.8, 1.45, 12.0, 5.5, "\n\n".join([f"[{part_idx + i + 1}] {source}" for i, source in enumerate(chunk)]), 11, GRAY)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    build_deck(Path("docs/reports/medical-insurance-top5-public-proxy.pptx"))
